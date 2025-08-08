# Universal Document Parser Module
# 
# This module requires Tesseract OCR to be installed on the system.
# Default Windows installation path:
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
# 
# For other operating systems, adjust the path accordingly:
# - Linux: usually '/usr/bin/tesseract'
# - macOS: usually '/usr/local/bin/tesseract' or '/opt/homebrew/bin/tesseract'

import re
import time
import io
from dataclasses import dataclass
from urllib.parse import urlparse
import requests
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from langchain.text_splitter import RecursiveCharacterTextSplitter
from typing import Generator
import docx
from pptx import Presentation
import pandas as pd
import eml_parser
import os

# Configure Tesseract path (uncomment and adjust for your system)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# File size limits (in MB)
MAX_FILE_SIZE_MB = 500
MAX_DOWNLOAD_TIME = 30

# Supported file extensions
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.eml', 
                       '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}

UNSUPPORTED_EXTENSIONS = {'.zip', '.rar', '.7z', '.tar', '.gz', '.bin', '.exe', '.dll', 
                         '.iso', '.img', '.dmg', '.deb', '.rpm'}


class UnsupportedDocumentError(Exception):
    """Custom exception for unsupported document types."""
    pass


class DocumentProcessingError(Exception):
    """Custom exception for document processing failures."""
    pass


@dataclass
class Chunk:
    """
    Data class to hold structured information about a text chunk.
    
    Attributes:
        text (str): The actual text content of the chunk
        page_number (int): The page number where this chunk originates
        source_label (str): A descriptive label indicating the source location and context
    """
    text: str
    page_number: int
    source_label: str


def _validate_file_extension(document_url: str) -> str:
    """
    Validate file extension before processing.
    
    Args:
        document_url (str): URL of the document
        
    Returns:
        str: Validated file extension
        
    Raises:
        UnsupportedDocumentError: If file extension is not supported
    """
    parsed_url = urlparse(document_url)
    path = parsed_url.path.lower()
    
    # Handle URLs with query parameters
    path = path.split('?')[0]
    
    if not path or '.' not in path:
        raise UnsupportedDocumentError("Unable to determine file type from URL")
    
    file_extension = os.path.splitext(path)[-1]
    
    if file_extension in UNSUPPORTED_EXTENSIONS:
        raise UnsupportedDocumentError(f"File type {file_extension} is not supported")
    
    if file_extension not in SUPPORTED_EXTENSIONS:
        raise UnsupportedDocumentError(f"File type {file_extension} is not supported")
    
    return file_extension


def _check_file_size(document_url: str) -> int:
    """
    Check file size before download to prevent memory issues.
    
    Args:
        document_url (str): URL of the document
        
    Returns:
        int: File size in bytes
        
    Raises:
        DocumentProcessingError: If file is too large or inaccessible
    """
    try:
        response = requests.head(document_url, timeout=10, allow_redirects=True)
        response.raise_for_status()
        
        if 'content-length' in response.headers:
            size_bytes = int(response.headers['content-length'])
            size_mb = size_bytes / (1024 * 1024)
            
            if size_mb > MAX_FILE_SIZE_MB:
                raise DocumentProcessingError(f"File too large: {size_mb:.1f}MB (max {MAX_FILE_SIZE_MB}MB)")
            
            print(f"File size: {size_mb:.1f}MB")
            return size_bytes
        else:
            print("Unable to determine file size, proceeding with caution")
            return 0
            
    except requests.RequestException as e:
        print(f"Warning: Could not check file size: {e}")
        return 0


def _parse_eml(content: bytes) -> list[tuple[str, int]]:
    """
    Parse EML email file extracting the email body.
    
    Args:
        content (bytes): EML document content
        
    Returns:
        list[tuple[str, int]]: List of (text, page_number) tuples
    """
    try:
        # Parse email from bytes
        parser = eml_parser.EmlParser()
        parsed_email = parser.decode_email_bytes(content)
        
        print("Successfully loaded EML email")
        
        # Extract email body
        body_text = ""
        if 'body' in parsed_email:
            # Handle different body formats
            body_data = parsed_email['body']
            if isinstance(body_data, list):
                # Multiple body parts
                for body_part in body_data:
                    if isinstance(body_part, dict) and 'content' in body_part:
                        body_text += body_part['content'] + "\n"
            elif isinstance(body_data, dict) and 'content' in body_data:
                body_text = body_data['content']
            elif isinstance(body_data, str):
                body_text = body_data
        
        # Also extract subject and headers for context
        email_content = ""
        if 'header' in parsed_email:
            header = parsed_email['header']
            if 'subject' in header:
                email_content += f"Subject: {header['subject']}\n"
            if 'from' in header:
                email_content += f"From: {header['from']}\n"
            if 'to' in header:
                email_content += f"To: {header['to']}\n"
            email_content += "\n"
        
        email_content += body_text
        
        print(f"Extracted {len(email_content)} characters from EML")
        
        if not email_content.strip():
            raise DocumentProcessingError("No content found in EML file")
            
        return [(email_content, 1)]
        
    except Exception as e:
        print(f"Error parsing EML document: {e}")
        raise DocumentProcessingError(f"Failed to parse EML: {e}")


def _parse_image(content: bytes, start_time: float, timeout: float) -> list[tuple[str, str]]:
    """
    Parse image file using OCR to extract text with intelligent line-based positioning.
    
    Args:
        content (bytes): Image file content
        start_time (float): Processing start time for timeout tracking
        timeout (float): Maximum processing time allowed
        
    Returns:
        list[tuple[str, str]]: List of (text, identifier) tuples
    """
    try:
        # Check timeout before starting OCR
        elapsed_time = time.time() - start_time
        if elapsed_time > timeout:
            raise DocumentProcessingError("Timeout reached before image OCR processing")
        
        # Open image from bytes
        image = Image.open(io.BytesIO(content))
        print(f"Successfully loaded image: {image.format}, {image.size}")
        
        # Check if image is too large (might cause memory issues)
        max_pixels = 50_000_000  # 50 megapixels
        if image.size[0] * image.size[1] > max_pixels:
            print(f"Image is very large ({image.size}), resizing for OCR...")
            # Calculate resize ratio to stay under limit
            ratio = (max_pixels / (image.size[0] * image.size[1])) ** 0.5
            new_size = (int(image.size[0] * ratio), int(image.size[1] * ratio))
            image = image.resize(new_size, Image.Resampling.LANCZOS)
            print(f"Resized to: {image.size}")
        
        # First, try to get text with bounding box data for better positioning
        try:
            # Check timeout again before expensive OCR operation
            elapsed_time = time.time() - start_time
            if elapsed_time > timeout:
                raise DocumentProcessingError("Timeout reached during image OCR processing")
            
            # Use pytesseract to get OCR data with bounding boxes
            ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
            
            # Reconstruct text with line information
            lines = {}
            for i, text in enumerate(ocr_data['text']):
                if text.strip():  # Only process non-empty text
                    top = ocr_data['top'][i]
                    # Group text by approximate line (using y-coordinate)
                    line_group = top // 20  # Group within 20 pixels vertically
                    if line_group not in lines:
                        lines[line_group] = []
                    lines[line_group].append((ocr_data['left'][i], text))
            
            # Sort lines by y-coordinate and words by x-coordinate within each line
            sorted_lines = []
            for line_y in sorted(lines.keys()):
                # Sort words in this line by x-coordinate (left position)
                words_in_line = sorted(lines[line_y], key=lambda x: x[0])
                line_text = ' '.join([word[1] for word in words_in_line])
                if line_text.strip():
                    sorted_lines.append((line_text, line_y))
            
            if sorted_lines:
                # Create text blocks with approximate line numbers
                full_ocr_text = ""
                
                for i, (line_text, line_y) in enumerate(sorted_lines):
                    full_ocr_text += line_text + "\n"
                
                print(f"Successfully extracted {len(sorted_lines)} text lines with positioning data")
                return [(full_ocr_text.strip(), f"Image Text with {len(sorted_lines)} positioned lines")]
            
        except Exception as ocr_data_error:
            print(f"Advanced OCR positioning failed, falling back to basic OCR: {ocr_data_error}")
        
        # Fallback: Basic OCR without positioning
        ocr_text = pytesseract.image_to_string(image)
        
        if ocr_text.strip():
            print(f"Extracted {len(ocr_text)} characters from image via basic OCR")
            # Return with improved identifier that won't cause "Unknown" labels
            return [(ocr_text, "Image Text")]
        else:
            print("No text found in image")
            raise DocumentProcessingError("No text content found in image")
        
    except DocumentProcessingError:
        raise
    except Exception as e:
        print(f"Error parsing image document: {e}")
        raise DocumentProcessingError(f"Failed to parse image: {e}")


def parse_document(document_url: str, timeout: int = 120) -> Generator[Chunk, None, None]:
    """
    Downloads and processes a document from a URL, supporting multiple file formats
    including PDF, DOCX, XLSX, PPTX, EML, and common image formats.
    
    Args:
        document_url (str): URL of the document to process
        timeout (int): Maximum time in seconds for total processing (default: 120)
        
    Returns:
        Generator[Chunk, None, None]: Generator yielding Chunk objects with text content and metadata
        
    Raises:
        UnsupportedDocumentError: If the file type is not supported
        DocumentProcessingError: If document processing fails
    """
    # Record start time for timeout tracking
    start_time = time.time()
    
    # Step 1: Validate file extension (5% of total time)
    try:
        file_extension = _validate_file_extension(document_url)
        print(f"Validated file extension: {file_extension}")
    except UnsupportedDocumentError:
        raise
    
    # Step 2: Check file size (5% of total time)
    try:
        _check_file_size(document_url)
    except DocumentProcessingError as e:
        print(f"File size check failed: {e}")
        raise UnsupportedDocumentError(str(e))
    
    # Step 3: Download document (20% of total time)
    download_timeout = min(MAX_DOWNLOAD_TIME, timeout * 0.2)
    try:
        print(f"Downloading document from: {document_url}")
        response = requests.get(document_url, timeout=download_timeout, stream=True)
        response.raise_for_status()
        
        # Read content with size limit
        content = b""
        max_size = MAX_FILE_SIZE_MB * 1024 * 1024
        for chunk in response.iter_content(chunk_size=8192):
            content += chunk
            if len(content) > max_size:
                raise DocumentProcessingError(f"File too large during download (>{MAX_FILE_SIZE_MB}MB)")
        
        print(f"Downloaded {len(content) / (1024*1024):.1f}MB in {time.time() - start_time:.1f}s")
        
    except requests.RequestException as e:
        raise DocumentProcessingError(f"Error downloading document: {e}")
    
    # Step 4: Process document using full timeout budget
    processing_start = time.time()
    
    try:
        if file_extension in ['.pdf']:
            yield from _parse_pdf(content, processing_start, timeout)
        elif file_extension in ['.docx', '.doc']:
            page_texts = _parse_docx(content)
            yield from _create_chunks_from_page_texts(page_texts)
        elif file_extension in ['.xlsx', '.xls']:
            page_texts = _parse_xlsx(content, file_extension)
            yield from _create_chunks_from_page_texts(page_texts)
        elif file_extension in ['.pptx', '.ppt']:
            page_texts = _parse_pptx(content, file_extension)
            yield from _create_chunks_from_page_texts(page_texts)
        elif file_extension in ['.eml']:
            page_texts = _parse_eml(content)
            yield from _create_chunks_from_page_texts(page_texts)
        elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
            page_texts = _parse_image(content, processing_start, timeout)
            yield from _create_chunks_from_page_texts(page_texts)
        else:
            raise UnsupportedDocumentError(f"File type {file_extension} is not supported")
            
    except UnsupportedDocumentError:
        raise
    except Exception as e:
        print(f"Error parsing document: {e}")
        raise DocumentProcessingError(f"Failed to process document: {str(e)}")


def _find_best_source_label(chunk_text: str, page_blocks: list, page_num: int, last_clause: str) -> tuple[str, str]:
    """
    Finds the best source label, preferring clauses but falling back to a line-number proxy.
    
    Args:
        chunk_text (str): The text content of the chunk
        page_blocks (list): List of text blocks from PyMuPDF page.get_text("blocks")
        page_num (int): Page number
        last_clause (str): The last known clause from previous chunks
        
    Returns:
        tuple[str, str]: (source_label, updated_last_clause)
    """
    clause_pattern = re.compile(r'^\s*([IVX\d]+\.[\d\.]*|[A-Z]\.[\d\.]*|\(\d+\)|[a-z]\))\s+.*', re.MULTILINE)

    chunk_start_y = -1
    chunk_block_text = ""
    
    # Find the first block that contains the start of our chunk text
    for block in page_blocks:
        if len(block) > 4:  # Ensure block has text content
            block_text = block[4]
            # Use a more flexible matching approach
            chunk_start = chunk_text.strip()[:100]  # First 100 chars for matching
            if chunk_start in block_text or block_text.strip() in chunk_text:
                chunk_start_y = block[1]  # The y0 coordinate
                chunk_block_text = block_text
                break
    
    # First, check if the chunk itself starts with a clause heading
    match = clause_pattern.match(chunk_text.strip())
    if match:
        clause = match.group(0).strip().split('\n')[0]
        return f"Page {page_num}, Clause: {clause}", clause
    
    # Second, check if the block containing this chunk starts with a clause
    if chunk_block_text:
        match = clause_pattern.match(chunk_block_text.strip())
        if match:
            clause = match.group(0).strip().split('\n')[0]
            return f"Page {page_num}, Clause: {clause}", clause
    
    # Third, search backwards through blocks to find the most recent clause
    if chunk_start_y != -1:
        for block in reversed(page_blocks):
            if len(block) > 4 and block[1] <= chunk_start_y:  # Block is above our chunk
                block_text = block[4]
                match = clause_pattern.search(block_text)
                if match:
                    clause = match.group(0).strip().split('\n')[0]
                    return f"Page {page_num}, Clause: {clause}", clause
    
    # Fourth, if we have a known clause from previous chunks on this page
    if last_clause and last_clause != "N/A":
        return f"Page {page_num}, Clause: {last_clause}", last_clause
    
    # Final fallback: use approximate line number based on y-coordinate
    if chunk_start_y != -1:
        approx_line = int(chunk_start_y / 15)  # Approximate line number (assuming 15pt line height)
        return f"Page {page_num}, Approx. Line: {approx_line}", last_clause
    
    # Ultimate fallback
    return f"Page {page_num}", last_clause


def _create_chunks_from_page_texts(page_texts: list[tuple[str, int]]) -> Generator[Chunk, None, None]:
    """
    Create chunks from page texts for non-PDF documents using traditional method.
    Enhanced to handle special page identifiers for XLSX, images, and provide intelligent source labels.
    
    Args:
        page_texts (list[tuple[str, int]]): List of (text, page_number_or_identifier) tuples
        
    Yields:
        Chunk: Individual chunk objects with metadata
    """
    if not page_texts:
        print("No text was extracted from the document")
        return
    
    # Join all page texts into a single document and create page map
    full_text = ""
    page_map = {}  # Maps character position to page number/identifier
    current_pos = 0
    
    for text, page_identifier in page_texts:
        page_map[current_pos] = page_identifier
        full_text += text + "\n\n"
        current_pos = len(full_text)
    
    print(f"Extracted total text length: {len(full_text)} characters")
    
    if not full_text.strip():
        print("No meaningful text content found")
        return
    
    # Define regex pattern for clause/section headings
    clause_pattern = re.compile(r'^\s*([IVX\d]+\.[\d\.]*|[A-Z]\.[\d\.]*|\(\d+\)|[a-z]\))\s+.*', re.MULTILINE)
    
    # Configure text splitter for chunking
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=150,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    # Split the full text into chunks
    text_chunks = splitter.split_text(full_text)
    print(f"Generated {len(text_chunks)} text chunks")
    
    # Create Chunk objects with metadata
    for chunk_text in text_chunks:
        if not chunk_text.strip():
            continue
            
        # Find the position of this chunk in the full text
        chunk_start = full_text.find(chunk_text)
        
        # Find the page identifier by looking up the chunk position in page_map
        page_identifier = 1  # Default to page 1
        for pos, page_id in sorted(page_map.items(), reverse=True):
            if chunk_start >= pos:
                page_identifier = page_id
                break
        
        # Handle special cases for different document types
        if isinstance(page_identifier, str):
            if "Image Text" in str(page_identifier):
                # Enhanced handling for image text with intelligent line positioning
                if "positioned lines" in str(page_identifier):
                    # Calculate approximate line number based on chunk position in text
                    lines_before = chunk_start // 50  # Rough estimate: 50 chars per line
                    source_label = f"Source: Image Text, Approx. Line: {lines_before + 1}"
                else:
                    # Fallback for basic image text
                    source_label = "Source: Image Text"
                page_number = 1
            elif "Spreadsheet Data" in str(page_identifier):
                # Handle XLSX/XLS sheet and row information
                source_label = f"Source: {page_identifier}"
                page_number = 1
            else:
                # Handle other special identifiers
                source_label = f"Source: {page_identifier}"
                page_number = 1
        else:
            # Regular page number handling
            page_number = page_identifier
            
            # Find the clause by searching backwards from chunk start
            clause_text = None
            if chunk_start >= 0:
                # Search backwards from chunk start to find the last clause match
                text_before_chunk = full_text[:chunk_start + len(chunk_text)]
                matches = list(clause_pattern.finditer(text_before_chunk))
                if matches:
                    # Get the last match (closest clause heading before this chunk)
                    last_match = matches[-1]
                    clause_line = last_match.group(0).strip()
                    # Extract just the first line of the match for cleaner display
                    clause_text = clause_line.split('\n')[0].strip()
            
            # Construct source label
            if clause_text:
                source_label = f"Page {page_number}, Clause: {clause_text}"
            else:
                # Use approximate line number as fallback
                approx_line = (chunk_start // 80) + 1  # Rough estimate based on character position
                source_label = f"Page {page_number}, Approx. Line: {approx_line}"
        
        # Create and yield Chunk object
        chunk_obj = Chunk(
            text=chunk_text,
            page_number=page_number if isinstance(page_identifier, int) else 1,
            source_label=source_label
        )
        yield chunk_obj


def _parse_pdf(content: bytes, start_time: float, timeout: float) -> Generator[Chunk, None, None]:
    """
    Parse PDF document using direct text extraction and OCR fallback with intelligent citation.
    
    Args:
        content (bytes): PDF document content
        start_time (float): Processing start time for timeout tracking
        timeout (float): Maximum processing time allowed
        
    Yields:
        Chunk: Individual chunk objects with intelligent source labels
    """
    try:
        # Load the PDF document from content
        doc = fitz.open(stream=io.BytesIO(content), filetype="pdf")
        total_pages = doc.page_count
        print(f"Successfully loaded PDF with {total_pages} pages")
        
    except Exception as e:
        print(f"Error loading PDF document: {e}")
        raise DocumentProcessingError(f"Failed to load PDF: {e}")
    
    # Configure text splitter for chunking
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=150,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    # Initialize counters
    pages_processed = 0
    timeout_reached = False
    chunks_generated = 0
    
    # Process each page in the document
    for page_num in range(total_pages):
        # Check timeout before processing each page
        elapsed_time = time.time() - start_time
        if elapsed_time > timeout:
            print(f"Timeout reached after {elapsed_time:.1f} seconds. Processed {pages_processed}/{total_pages} pages.")
            timeout_reached = True
            break
            
        try:
            page = doc[page_num]
            print(f"Processing page {page_num + 1}/{total_pages} (elapsed: {elapsed_time:.1f}s)")
            
            # Extract structured text blocks with positional data
            page_blocks = page.get_text("blocks")
            
            # Check if page has meaningful content
            full_page_text = "\n".join([block[4] for block in page_blocks if len(block) > 4 and block[4].strip()])
            
            if not full_page_text.strip():
                print(f"  Page {page_num + 1}: Skipping empty page")
                pages_processed += 1
                continue
            
            # Heuristic: if extracted text is less than 100 characters after stripping,
            # assume it's a scanned page or contains mostly images/diagrams
            if len(full_page_text.strip()) < 100:
                # This appears to be a scanned page - use OCR
                print(f"  Page {page_num + 1}: Using OCR (scanned content detected)")
                
                # Check timeout again before expensive OCR operation
                elapsed_time = time.time() - start_time
                if elapsed_time > timeout:
                    print(f"Timeout reached before OCR processing on page {page_num + 1}")
                    timeout_reached = True
                    break
                
                # Render page as high-resolution image
                pixmap = page.get_pixmap(dpi=300)
                
                # Convert pixmap to PIL Image for pytesseract
                img_data = pixmap.tobytes("ppm")
                pil_image = Image.open(io.BytesIO(img_data))
                
                # Perform OCR on the image
                ocr_text = pytesseract.image_to_string(pil_image)
                
                # For OCR content, create simple text blocks structure
                if ocr_text.strip():
                    # Create artificial blocks for OCR text with estimated coordinates
                    ocr_blocks = []
                    lines = ocr_text.split('\n')
                    for i, line in enumerate(lines):
                        if line.strip():
                            # Create artificial block with estimated y-coordinate
                            y_coord = i * 15  # Assume 15pt line spacing
                            ocr_blocks.append((0, y_coord, 0, y_coord + 15, line))
                    page_blocks = ocr_blocks
                    full_page_text = ocr_text
            else:
                # Digital text found - use direct extraction
                print(f"  Page {page_num + 1}: Using direct text extraction")
            
            # Chunk the page text
            page_chunks_text = splitter.split_text(full_page_text)
            
            # Track the last known clause for this page
            last_clause = "N/A"
            
            # Generate chunks with intelligent source labels
            for chunk_text in page_chunks_text:
                if not chunk_text.strip():
                    continue
                    
                source_label, last_clause = _find_best_source_label(
                    chunk_text, page_blocks, page_num + 1, last_clause
                )
                
                chunk_obj = Chunk(
                    text=chunk_text,
                    page_number=page_num + 1,
                    source_label=source_label
                )
                yield chunk_obj
                chunks_generated += 1
            
            pages_processed += 1
                
        except Exception as e:
            print(f"Error processing page {page_num + 1}: {e}")
            # Continue with next page rather than failing completely
            pages_processed += 1
            continue
    
    # Close the document to free memory
    doc.close()
    
    # Report processing results
    total_elapsed = time.time() - start_time
    if timeout_reached:
        print(f"PDF processing stopped due to timeout. Processed {pages_processed}/{total_pages} pages, generated {chunks_generated} chunks in {total_elapsed:.1f}s")
    else:
        print(f"PDF processing completed. Processed {pages_processed}/{total_pages} pages, generated {chunks_generated} chunks in {total_elapsed:.1f}s")
    
    if chunks_generated == 0:
        raise DocumentProcessingError("No content could be extracted from the PDF")


def _parse_docx(content: bytes) -> list[tuple[str, int]]:
    """
    Parse DOCX/DOC document extracting text from paragraphs and tables.
    
    Args:
        content (bytes): DOCX document content
        
    Returns:
        list[tuple[str, int]]: List of (text, page_number) tuples
    """
    try:
        # Open document from bytes
        doc = docx.Document(io.BytesIO(content))
        print("Successfully loaded DOCX document")
        
        extracted_text = ""
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                extracted_text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    extracted_text += " | ".join(row_text) + "\n"
        
        print(f"Extracted {len(extracted_text)} characters from DOCX")
        
        if not extracted_text.strip():
            raise DocumentProcessingError("No text content found in DOCX document")
            
        return [(extracted_text, 1)]
        
    except Exception as e:
        print(f"Error parsing DOCX document: {e}")
        raise DocumentProcessingError(f"Failed to parse DOCX: {e}")


def _parse_xlsx(content: bytes, file_extension: str) -> list[tuple[str, str]]:
    """
    Parse XLSX/XLS document extracting structured data from all sheets with enhanced location tracking.
    
    Args:
        content (bytes): XLSX document content
        file_extension (str): File extension (.xlsx or .xls)
        
    Returns:
        list[tuple[str, str]]: List of (text, location_identifier) tuples
    """
    try:
        # Handle both .xlsx and .xls files
        if file_extension == '.xls':
            # For .xls files, try to read with xlrd engine
            xls = pd.ExcelFile(io.BytesIO(content), engine='xlrd')
        else:
            # For .xlsx files, use default openpyxl engine
            xls = pd.ExcelFile(io.BytesIO(content))
            
        print(f"Successfully loaded {file_extension.upper()} document with sheets: {xls.sheet_names}")
        
        all_text_parts = []
        
        # Process each sheet
        for sheet_name in xls.sheet_names:
            try:
                if file_extension == '.xls':
                    df = pd.read_excel(xls, sheet_name=sheet_name, engine='xlrd')
                else:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    
                sheet_text = f"Sheet: {sheet_name}\n"
                
                # Handle empty sheets
                if df.empty:
                    print(f"  Sheet '{sheet_name}' is empty, skipping")
                    continue
                
                # Convert each row to a readable string with enhanced location tracking
                for index, row in df.iterrows():
                    row_parts = []
                    for col_name, value in row.items():
                        if pd.notna(value):  # Skip NaN values
                            row_parts.append(f"{col_name}: {value}")
                    
                    if row_parts:
                        # Enhanced row formatting with sheet and row information
                        row_text = f"Sheet '{sheet_name}', Row {index + 2}: {', '.join(row_parts)}\n"
                        sheet_text += row_text
                
                if len(sheet_text) > len(f"Sheet: {sheet_name}\n"):  # Has content beyond header
                    all_text_parts.append(sheet_text)
                else:
                    print(f"  Sheet '{sheet_name}' has no meaningful data, skipping")
                
            except Exception as e:
                print(f"Error processing sheet '{sheet_name}': {e}")
                continue
        
        if not all_text_parts:
            raise DocumentProcessingError("No readable data found in spreadsheet")
        
        combined_text = "\n\n".join(all_text_parts)
        print(f"Extracted {len(combined_text)} characters from {file_extension.upper()}")
        
        # Return with special identifier for spreadsheet content
        return [(combined_text, f"{file_extension.upper()[1:]} Spreadsheet Data")]
        
    except Exception as e:
        print(f"Error parsing {file_extension.upper()} document: {e}")
        raise DocumentProcessingError(f"Failed to parse {file_extension.upper()}: {e}")


def _parse_pptx(content: bytes, file_extension: str) -> list[tuple[str, int]]:
    """
    Comprehensive PPTX/PPT parser that extracts text from slides, notes, and tables.
    
    Args:
        content (bytes): PPTX document content
        file_extension (str): File extension (.pptx or .ppt)
        
    Returns:
        list[tuple[str, int]]: List of (text, page_number) tuples
    """
    try:
        prs = Presentation(io.BytesIO(content))
        page_texts = []
        
        print(f"Successfully loaded {file_extension.upper()} presentation with {len(prs.slides)} slides")
        
        if len(prs.slides) == 0:
            raise DocumentProcessingError("Presentation contains no slides")
        
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            
            try:
                # Extract notes text if available
                if hasattr(slide, 'notes_slide') and slide.notes_slide and hasattr(slide.notes_slide, 'notes_text_frame'):
                    if slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text.strip():
                            slide_text += f"Notes: {notes_text}\n"
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    try:
                        # Extract text from text frames
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                            shape_text = shape.text_frame.text
                            if shape_text.strip():
                                slide_text += shape_text + "\n"
                        
                        # Extract text from tables
                        if hasattr(shape, 'has_table') and shape.has_table:
                            for row in shape.table.rows:
                                row_text = ""
                                for cell in row.cells:
                                    if hasattr(cell, 'text_frame'):
                                        cell_text = cell.text_frame.text
                                        if cell_text.strip():
                                            row_text += cell_text + "\t"
                                if row_text.strip():
                                    slide_text += row_text + "\n"
                    except Exception as shape_error:
                        print(f"  Warning: Error processing shape on slide {i+1}: {shape_error}")
                        continue
                
                # Only add slide if it contains meaningful text
                if slide_text.strip():
                    page_texts.append((slide_text, i + 1))
                else:
                    print(f"  Slide {i+1}: No text content found")
                    
            except Exception as slide_error:
                print(f"  Error processing slide {i+1}: {slide_error}")
                # Add a placeholder for the problematic slide
                page_texts.append((f"Error processing slide content: {str(slide_error)}", i + 1))
                continue
        
        print(f"Extracted text from {len(page_texts)} slides with content")
        
        if not page_texts:
            raise DocumentProcessingError("No text content found in any slides")
            
        return page_texts
        
    except Exception as e:
        print(f"Error parsing {file_extension.upper()} document: {e}")
        raise DocumentProcessingError(f"Failed to parse {file_extension.upper()}: {e}")